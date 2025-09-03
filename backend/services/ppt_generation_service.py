#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AI智能学习系统 - 业务服务 - ppt_generation_service.py

Description:
    PPT生成服务，提供基于AI的PPT自动生成功能。

Author: Chang Xinglong
Date: 2025-01-21
Version: 1.0.0
License: Apache License 2.0
"""

import os
import json
import uuid
from datetime import datetime
from typing import Dict, List, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
# 移除有问题的导入
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

from utils.database import db
from models.document import Document
from models.ppt_template import PPTTemplate
from services.llm_service import llm_service
from services.document_service import get_document_service
from utils.logger import get_logger

logger = get_logger(__name__)

class PPTGenerationService:
    """
    PPT生成服务类
    """

    def __init__(self):
        self.upload_folder = 'uploads/ppt'
        os.makedirs(self.upload_folder, exist_ok=True)

        # PPT模板配置
        self.templates = {
            'default': {
                'title_font_size': 44,
                'content_font_size': 24,
                'title_color': RGBColor(31, 73, 125),
                'content_color': RGBColor(68, 68, 68),
                'background_color': RGBColor(255, 255, 255)
            },
            'education': {
                'title_font_size': 40,
                'content_font_size': 22,
                'title_color': RGBColor(46, 125, 50),
                'content_color': RGBColor(33, 33, 33),
                'background_color': RGBColor(248, 249, 250)
            }
        }

    def generate_ppt_from_text(self, user_id: str, tenant_id: str, content: str,
                              title: str = "AI生成的演示文稿",
                              template: str = "default",
                              template_id: Optional[str] = None) -> Dict[str, Any]:
        """
        根据文本内容生成PPT

        Args:
            user_id: 用户ID
            tenant_id: 租户ID
            content: 文本内容
            title: PPT标题
            template: 模板类型

        Returns:
            Dict: 生成结果
        """
        try:
            # 使用AI分析内容并生成PPT结构
            ppt_structure = self._analyze_content_for_ppt(content, title)

            if not ppt_structure['success']:
                return {
                    'success': False,
                    'error': ppt_structure.get('error', 'PPT结构分析失败')
                }

            # 生成PPT文件
            ppt_result = self._create_ppt_file(
                ppt_structure['data'],
                title,
                template,
                template_id,
                tenant_id
            )

            if not ppt_result['success']:
                return ppt_result

            # 保存到文档管理系统
            document_result = self._save_ppt_to_documents(
                user_id,
                tenant_id,
                ppt_result['file_path'],
                title,
                content
            )

            if not document_result['success']:
                return document_result

            return {
                'success': True,
                'data': {
                    'document_id': document_result['document_id'],
                    'file_path': ppt_result['file_path'],
                    'download_url': f'http://localhost:5001/api/document/{document_result["document_id"]}/download',
                    'title': title,
                    'slides_count': len(ppt_structure['data']['slides']),
                    'created_at': datetime.now().isoformat()
                }
            }

        except Exception as e:
            logger.error(f"PPT生成失败: {str(e)}")
            return {
                'success': False,
                'error': f'PPT生成失败: {str(e)}'
            }

    def _analyze_content_for_ppt(self, content: str, title: str) -> Dict[str, Any]:
        """
        分析内容并生成PPT结构

        Args:
            content: 文本内容
            title: PPT标题

        Returns:
            Dict: 分析结果
        """
        try:
            prompt = f"""
你是一个专业的PPT制作专家。请根据以下内容，设计一个结构清晰、内容丰富的PPT演示文稿。

标题：{title}
内容：{content}

请按照以下JSON格式返回PPT结构：
{{
    "title": "PPT主标题",
    "slides": [
        {{
            "type": "title",
            "title": "幻灯片标题",
            "subtitle": "副标题（可选）"
        }},
        {{
            "type": "content",
            "title": "幻灯片标题",
            "content": [
                "要点1",
                "要点2",
                "要点3"
            ]
        }},
        {{
            "type": "two_column",
            "title": "幻灯片标题",
            "left_content": [
                "左侧要点1",
                "左侧要点2"
            ],
            "right_content": [
                "右侧要点1",
                "右侧要点2"
            ]
        }}
    ]
}}

要求：
1. 第一张幻灯片必须是标题页（type: "title"）
2. 内容要点要简洁明了，每个要点不超过50字
3. 合理安排幻灯片数量，建议5-15张
4. 根据内容特点选择合适的幻灯片类型
5. 确保逻辑清晰，层次分明

请只返回JSON格式的结果，不要包含其他文字。
"""

            result = llm_service.generate_text(
                prompt=prompt,
                temperature=0.3
            )

            if not result:
                return {
                    'success': False,
                    'error': 'AI分析内容失败'
                }

            # 记录LLM返回的原始结果用于调试
            logger.info(f"LLM返回的原始结果: {result[:500]}...")  # 只记录前500个字符

            # 清理结果，移除可能的markdown代码块标记
            cleaned_result = result.strip()
            if cleaned_result.startswith('```json'):
                cleaned_result = cleaned_result[7:]
            if cleaned_result.startswith('```'):
                cleaned_result = cleaned_result[3:]
            if cleaned_result.endswith('```'):
                cleaned_result = cleaned_result[:-3]
            cleaned_result = cleaned_result.strip()

            # 解析JSON结果
            try:
                ppt_data = json.loads(cleaned_result)
                return {
                    'success': True,
                    'data': ppt_data
                }
            except json.JSONDecodeError as e:
                logger.error(f"PPT结构JSON解析失败: {str(e)}")
                logger.error(f"清理后的结果: {cleaned_result[:200]}...")  # 记录清理后的结果
                return {
                    'success': False,
                    'error': f'PPT结构解析失败: {str(e)}'
                }

        except Exception as e:
            logger.error(f"PPT内容分析失败: {str(e)}")
            return {
                'success': False,
                'error': f'内容分析失败: {str(e)}'
            }

    def _get_template_config(self, template_id: Optional[str], template: str, tenant_id: str) -> Dict[str, Any]:
        """
        获取模板配置

        Args:
            template_id: 模板ID（优先使用）
            template: 模板类型名称
            tenant_id: 租户ID

        Returns:
            Dict: 模板配置
        """
        try:
            # 如果提供了template_id，优先从数据库获取
            if template_id:
                db_template = PPTTemplate.query.filter_by(
                    id=template_id,
                    tenant_id=tenant_id,
                    is_active=True
                ).first()

                if db_template and db_template.config:
                    return db_template.config

            # 如果没有找到数据库模板，使用内置模板
            return self.templates.get(template, self.templates['default'])

        except Exception as e:
            logger.warning(f"获取模板配置失败，使用默认配置: {str(e)}")
            return self.templates['default']
    
    def _parse_color(self, color_value):
        """解析颜色值，支持字符串和RGBColor对象"""
        if isinstance(color_value, str):
            # 移除#号并转换为RGB
            hex_color = color_value.lstrip('#')
            if len(hex_color) == 6:
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                return RGBColor(r, g, b)
        elif hasattr(color_value, 'rgb'):
            return color_value
        # 默认返回黑色
        return RGBColor(0, 0, 0)
    
    def _load_template_file(self, template_id: Optional[str], tenant_id: str):
        """
        加载模板文件
        
        Args:
            template_id: 模板ID
            tenant_id: 租户ID
            
        Returns:
            Presentation: PPT演示文稿对象
        """
        try:
            # 如果提供了template_id，尝试从数据库获取模板文件
            if template_id:
                db_template = PPTTemplate.query.filter_by(
                    id=template_id,
                    tenant_id=tenant_id,
                    is_active=True
                ).first()
                
                if db_template and db_template.template_file_path:
                    template_path = db_template.template_file_path
                    if os.path.exists(template_path):
                        logger.info(f"使用用户模板文件: {template_path}")
                        return Presentation(template_path)
                    else:
                        logger.warning(f"模板文件不存在: {template_path}")
            
            # 如果没有找到用户模板文件，使用空白演示文稿
            logger.info("使用空白演示文稿")
            return Presentation()
            
        except Exception as e:
            logger.warning(f"加载模板文件失败，使用空白演示文稿: {str(e)}")
            return Presentation()

    def _create_ppt_file(self, ppt_data: Dict[str, Any], title: str, template: str,
                        template_id: Optional[str] = None, tenant_id: str = "default") -> Dict[str, Any]:
        """
        创建PPT文件

        Args:
            ppt_data: PPT数据结构
            title: PPT标题
            template: 模板类型

        Returns:
            Dict: 创建结果
        """
        try:
            # 尝试加载用户上传的模板文件
            prs = self._load_template_file(template_id, tenant_id)
            
            # 获取模板配置
            template_config = self._get_template_config(template_id, template, tenant_id)
            
            # 如果使用了用户模板文件，清空现有幻灯片（保留布局）
            if template_id and len(prs.slides) > 0:
                # 记录可用的布局
                available_layouts = prs.slide_layouts
                # 删除所有现有幻灯片
                try:
                    while len(prs.slides) > 0:
                        # 获取第一个幻灯片的关系ID
                        slide_part = prs.slides[0].part
                        # 从幻灯片列表中移除
                        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
                        # 从部件中删除关系
                        if hasattr(slide_part, 'package'):
                            prs.part.drop_rel(slide_part.partname)
                except Exception as clear_error:
                    logger.warning(f"清空模板幻灯片时出现问题，继续使用原模板: {str(clear_error)}")
                    # 如果清空失败，重新加载模板
                    prs = self._load_template_file(template_id, tenant_id)

            # 创建幻灯片
            for slide_data in ppt_data.get('slides', []):
                slide_type = slide_data.get('type', 'content')

                if slide_type == 'title':
                    self._create_title_slide(prs, slide_data, template_config)
                elif slide_type == 'content':
                    self._create_content_slide(prs, slide_data, template_config)
                elif slide_type == 'two_column':
                    self._create_two_column_slide(prs, slide_data, template_config)
                else:
                    # 默认创建内容幻灯片
                    self._create_content_slide(prs, slide_data, template_config)

            # 生成文件名
            filename = f"{uuid.uuid4().hex}_{title.replace(' ', '_')}.pptx"
            file_path = os.path.join(self.upload_folder, filename)

            # 保存文件
            prs.save(file_path)

            return {
                'success': True,
                'file_path': file_path,
                'filename': filename
            }

        except Exception as e:
            logger.error(f"PPT文件创建失败: {str(e)}")
            return {
                'success': False,
                'error': f'PPT文件创建失败: {str(e)}'
            }

    def _create_title_slide(self, prs, slide_data: Dict, template_config: Dict):
        """
        创建标题幻灯片
        """
        slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
        slide = prs.slides.add_slide(slide_layout)

        # 设置标题
        title = slide.shapes.title
        title.text = slide_data.get('title', '')
        title.text_frame.paragraphs[0].font.size = Pt(template_config['title_font_size'])
        title.text_frame.paragraphs[0].font.color.rgb = self._parse_color(template_config['title_color'])
        # 设置居中对齐 - 使用数字常量 1 表示居中
        title.text_frame.paragraphs[0].alignment = 1

        # 设置副标题
        if 'subtitle' in slide_data and slide.placeholders[1]:
            subtitle = slide.placeholders[1]
            subtitle.text = slide_data['subtitle']
            subtitle.text_frame.paragraphs[0].font.size = Pt(template_config['content_font_size'])
            subtitle.text_frame.paragraphs[0].font.color.rgb = self._parse_color(template_config['content_color'])
            # 设置居中对齐 - 使用数字常量 1 表示居中
            subtitle.text_frame.paragraphs[0].alignment = 1

    def _create_content_slide(self, prs, slide_data: Dict, template_config: Dict):
        """
        创建内容幻灯片
        """
        slide_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(slide_layout)

        # 设置标题
        title = slide.shapes.title
        title.text = slide_data.get('title', '')
        title.text_frame.paragraphs[0].font.size = Pt(template_config['title_font_size'] - 8)
        title.text_frame.paragraphs[0].font.color.rgb = self._parse_color(template_config['title_color'])

        # 设置内容
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()

        content_list = slide_data.get('content', [])
        for i, content_item in enumerate(content_list):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = content_item
            p.font.size = Pt(template_config['content_font_size'])
            p.font.color.rgb = self._parse_color(template_config['content_color'])
            p.level = 0

    def _create_two_column_slide(self, prs, slide_data: Dict, template_config: Dict):
        """
        创建双栏幻灯片
        """
        slide_layout = prs.slide_layouts[3]  # 两个内容布局
        slide = prs.slides.add_slide(slide_layout)

        # 设置标题
        title = slide.shapes.title
        title.text = slide_data.get('title', '')
        title.text_frame.paragraphs[0].font.size = Pt(template_config['title_font_size'] - 8)
        title.text_frame.paragraphs[0].font.color.rgb = self._parse_color(template_config['title_color'])

        # 左侧内容
        left_placeholder = slide.placeholders[1]
        left_text_frame = left_placeholder.text_frame
        left_text_frame.clear()

        left_content = slide_data.get('left_content', [])
        for i, content_item in enumerate(left_content):
            p = left_text_frame.paragraphs[0] if i == 0 else left_text_frame.add_paragraph()
            p.text = content_item
            p.font.size = Pt(template_config['content_font_size'])
            p.font.color.rgb = self._parse_color(template_config['content_color'])
            p.level = 0

        # 右侧内容
        right_placeholder = slide.placeholders[2]
        right_text_frame = right_placeholder.text_frame
        right_text_frame.clear()

        right_content = slide_data.get('right_content', [])
        for i, content_item in enumerate(right_content):
            p = right_text_frame.paragraphs[0] if i == 0 else right_text_frame.add_paragraph()
            p.text = content_item
            p.font.size = Pt(template_config['content_font_size'])
            p.font.color.rgb = self._parse_color(template_config['content_color'])
            p.level = 0

    def _save_ppt_to_documents(self, user_id: str, tenant_id: str, file_path: str,
                              title: str, original_content: str) -> Dict[str, Any]:
        """
        将PPT保存到文档管理系统

        Args:
            user_id: 用户ID
            tenant_id: 租户ID
            file_path: PPT文件路径
            title: PPT标题
            original_content: 原始内容

        Returns:
            Dict: 保存结果
        """
        try:
            document_service = get_document_service()

            # 创建文档记录
            document = Document(
                id=str(uuid.uuid4()),
                title=title,
                description=f"AI生成的PPT演示文稿，基于内容：{original_content[:100]}...",
                filename=os.path.basename(file_path),
                file_path=file_path,
                file_size=os.path.getsize(file_path),
                file_type='pptx',
                mime_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                file_hash=document_service.calculate_file_hash(file_path),
                user_id=user_id,
                tenant_id=tenant_id,
                parse_status='completed'
            )

            db.session.add(document)
            db.session.commit()

            return {
                'success': True,
                'document_id': document.id
            }

        except Exception as e:
            logger.error(f"PPT文档保存失败: {str(e)}")
            db.session.rollback()
            return {
                'success': False,
                'error': f'文档保存失败: {str(e)}'
            }

    def detect_ppt_request(self, message: str) -> Dict[str, Any]:
        """
        检测用户消息是否包含PPT生成请求

        Args:
            message: 用户消息

        Returns:
            Dict: 检测结果
        """
        ppt_keywords = [
            'ppt', 'PPT', 'powerpoint', 'PowerPoint', '演示文稿', '幻灯片',
            '制作ppt', '生成ppt', '创建ppt', '做个ppt', '帮我做ppt',
            '制作演示文稿', '生成演示文稿', '创建演示文稿'
        ]

        message_lower = message.lower()

        # 检查是否包含PPT相关关键词
        contains_ppt_keyword = any(keyword.lower() in message_lower for keyword in ppt_keywords)

        if contains_ppt_keyword:
            # 尝试提取主题
            topic = self._extract_ppt_topic(message)
            return {
                'is_ppt_request': True,
                'topic': topic,
                'confidence': 0.9 if topic else 0.7
            }

        return {
            'is_ppt_request': False,
            'topic': None,
            'confidence': 0.0
        }

    def _extract_ppt_topic(self, message: str) -> Optional[str]:
        """
        从消息中提取PPT主题

        Args:
            message: 用户消息

        Returns:
            Optional[str]: 提取的主题
        """
        # 简单的主题提取逻辑
        # 可以后续使用更复杂的NLP技术

        # 常见的主题引导词
        topic_indicators = ['关于', '主题是', '内容是', '讲', '介绍', '分析']

        for indicator in topic_indicators:
            if indicator in message:
                parts = message.split(indicator)
                if len(parts) > 1:
                    topic = parts[1].strip()
                    # 清理主题文本
                    topic = topic.replace('的ppt', '').replace('的PPT', '').replace('的演示文稿', '')
                    if topic:
                        return topic[:50]  # 限制主题长度

        return None

# 创建服务实例
ppt_generation_service = PPTGenerationService()
